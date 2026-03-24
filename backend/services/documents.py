"""
backend/services/documents.py — Document text extraction.

Provides extract_document_text() which accepts an uploaded file path and
returns the extracted slide/page list in the unified format used by
/upload-document, /generate-study-materials, and /generate-quiz.
"""
from __future__ import annotations

import logging
import os

logger = logging.getLogger(__name__)


def extract_slides_from_file(temp_path: str, safe_name: str) -> list[dict]:
    """
    Extract text from a PDF, DOCX, or PPTX file.

    Returns a list of dicts:
        [{'slide_number': int, 'title': str, 'content': [str, ...], 'notes': str}, ...]

    Raises ValueError if the file type is unsupported.
    Raises RuntimeError if extraction yields no readable text.
    """
    filename = safe_name.lower()
    extracted_slides: list[dict] = []

    if filename.endswith(('.pptx', '.ppt')):
        extracted_slides = _extract_pptx(temp_path)
    elif filename.endswith('.pdf'):
        extracted_slides = _extract_pdf(temp_path)
    elif filename.endswith('.docx'):
        extracted_slides = _extract_docx(temp_path, safe_name)
    else:
        raise ValueError(f'Unsupported file type: {filename}')

    total_text = " ".join(" ".join(s.get('content', [])) for s in extracted_slides)
    if len(total_text.strip()) < 30:
        raise RuntimeError(
            'Could not extract readable text. '
            'The file may be scanned/image-based or empty.'
        )

    return extracted_slides


def _extract_pptx(temp_path: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(temp_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        slide_title = f"Slide {i}"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                is_title = (
                    hasattr(shape, "is_placeholder") and shape.is_placeholder and
                    hasattr(shape, "placeholder_format") and shape.placeholder_format is not None and
                    shape.placeholder_format.idx == 0
                )
                if is_title:
                    slide_title = text
                else:
                    slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(row_text)
        notes_text = ""
        if slide.has_notes_slide:
            nf = slide.notes_slide.notes_text_frame
            if nf:
                notes_text = nf.text.strip()
        slides.append({
            'slide_number': i, 'title': slide_title,
            'content': slide_texts, 'notes': notes_text
        })
    return slides


def _extract_pdf(temp_path: str) -> list[dict]:
    import PyPDF2
    slides = []
    with open(temp_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            slides.append({
                'slide_number': i, 'title': f"Page {i}",
                'content': [text] if text.strip() else [], 'notes': ''
            })
    return slides


def _extract_docx(temp_path: str, safe_name: str) -> list[dict]:
    import docx as _docx
    doc = _docx.Document(temp_path)
    slides = []
    current_section = {"title": safe_name, "content": [], "slide_number": 1, "notes": ""}
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        if para.style.name.startswith('Heading'):
            if current_section["content"]:
                slides.append(current_section)
            current_section = {
                "title": para.text.strip(), "content": [],
                "slide_number": len(slides) + 1, "notes": ""
            }
        else:
            current_section["content"].append(para.text.strip())
    if current_section["content"]:
        slides.append(current_section)
    if not slides:
        all_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        slides.append({
            'slide_number': 1, 'title': safe_name,
            'content': [all_text], 'notes': ''
        })
    return slides
